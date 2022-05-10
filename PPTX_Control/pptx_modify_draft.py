from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

import openpyxl
import os

if __name__ == "__main__":

    CUR_PATH = os.getcwd()
    SRC_FILE = CUR_PATH + "\\ip_list_1.xlsx"

    wb = openpyxl.load_workbook(SRC_FILE)
    ws = wb.active

    for row in range(2, ws.max_row + 1):
        schNo = ws[f'A{row}'].value
        schOrg = ws[f'B{row}'].value
        schName = ws[f'D{row}'].value
        fileName = ws[f'I{row}'].value
        if str(ws[f'J{row}'].value) != "None":
            ip_teacher = ws[f'J{row}'].value.split(",")
        else:
            ip_teacher = ""

        if str(ws[f'K{row}'].value) != "None":
            ip_student = ws[f'K{row}'].value.split(",")
        else:
            ip_student = ""

        if str(ws[f'L{row}'].value) != "None":
            ip_wireless = ws[f'L{row}'].value.split(",")
        else:
            ip_wireless = ""

        if str(ws[f'N{row}'].value) != "None":
            ip_etc = ws[f'N{row}'].value.split(",")
        else:
            ip_etc = ""

        if len(ip_teacher) > 1:
            print(ip_teacher)

        prs = Presentation(".\\1.pptx")
        main_slide = prs.slides[0]

        for shape in prs.slides[0].shapes:
            if not shape.has_text_frame:
                continue
            if "기관명 :" in shape.text_frame.paragraphs[0].text:
                # print(shape.text_frame.paragraphs[0].text)
                shape.text_frame.paragraphs[0].text = f"기관명 : {schOrg}-{schName}"
                for run in shape.text_frame.paragraphs[0].runs:
                    run.font.size = Pt(16)
                    run.font.name = '맑은 고딕'
                    run.font.bold = True

        for shape in prs.slides[0].shapes:
            if not shape.has_table:
                continue
            if "Port" in shape.table.cell(0, 0).text:
                # print(shape.table.cell(0, 0).text)
                if ip_teacher != "":
                    ip_net = ip_teacher[0]
                    if len(ip_teacher) > 1:
                        ip_net += "\n" + ip_teacher[1]
                    shape.table.cell(2, 3).text = ip_net
                else:
                    shape.table.cell(2, 3).text = ""
                for paragraph in shape.table.cell(2, 3).text_frame.paragraphs:
                    paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(12)
                        run.font.name = '맑은 고딕'
                        run.font.bold = True

                if ip_student != "":
                    ip_net = ip_student[0]
                    if len(ip_student) > 1:
                        ip_net += "\n" + ip_student[1]
                    shape.table.cell(3, 3).text = ip_net
                else:
                    shape.table.cell(3, 3).text = ""
                for paragraph in shape.table.cell(3, 3).text_frame.paragraphs:
                    paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(12)
                        run.font.name = '맑은 고딕'
                        run.font.bold = True

                if ip_etc != "":
                    ip_net = ip_etc[0]
                    if len(ip_etc) > 1:
                        ip_net += "\n" + ip_etc[1]
                    shape.table.cell(4, 3).text = ip_net
                else:
                    shape.table.cell(4, 3).text = ""
                for paragraph in shape.table.cell(4, 3).text_frame.paragraphs:
                    paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(12)
                        run.font.name = '맑은 고딕'
                        run.font.bold = True

                if ip_wireless != "":
                    ip_net = ip_wireless[0]
                    if len(ip_wireless) > 1:
                        ip_net += "\n" + ip_wireless[1]
                    shape.table.cell(6, 3).text = ip_net
                else:
                    shape.table.cell(6, 3).text = ""
                for paragraph in shape.table.cell(6, 3).text_frame.paragraphs:
                    paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(12)
                        run.font.name = '맑은 고딕'
                        run.font.bold = True
        # if schNo < 10:
        #     fileNo = "00" + str(schNo)
        # elif schNo < 100:
        #     fileNo = "0" + str(schNo)
        # else:
        #     fileNo = str(schNo)
        prs.save(f".\\{fileName}_망구성도.pptx")
