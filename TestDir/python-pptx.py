from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

prs = Presentation("D:\\Desktop\\다. 대은초등학교_망구성도.pptx")

# text_runs will be populated with a list of strings,
# one for each text run in presentation

i=0
j=0
text_runs = []
print(prs.slides[0].shapes[2].table.cell(5,2).text)

prs.slides[0].shapes[2].table.cell(5,2).text = "기타(서비스망)"
prs.slides[0].shapes[2].table.cell(5,2).text_frame.paragraphs[0].font.name = '나눔고딕 Bold'
prs.slides[0].shapes[2].table.cell(5,2).text_frame.paragraphs[0].font.size = Pt(10)
prs.slides[0].shapes[2].table.cell(5,2).text_frame.paragraphs[0].font.bold = True
prs.slides[0].shapes[2].table.cell(5,2).text_frame.paragraphs[0].font.italic = None
prs.slides[0].shapes[2].table.cell(5,2).text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

print(prs.slides[0].shapes[2].table.cell(5,2).text)
prs.save("D:\\Desktop\\다. 대은초등학교_망구성도.pptx")

'''for slide in prs.slides:
    j = j + 1
    print(str(j) + " Slide")
    for shape in slide.shapes:
        if shape.has_table:
            print(shape.has_table)
            for r in shape.table.rows:
                print(r)
                for c in r.cells:
                    text_runs.append(c.text_frame.text)
                    #print(c.text_frame.text)
        i = i + 1
        print(str(i) + " Shape")
  #  print(text_runs)
    text_runs = []

#print(prs.slides[0].shapes[6].table.rows.cells(1,0).text)'''
'''
for slide in prs.slides:
    i = i + 1
    print(str(i) + "Slide")
    for shape in slide.shapes:
        j = j + 1
        print(str(j) + "Shape")
        print(shape.has_table)
print(prs.slides[1].shapes[3].table.cell(1,1).text)
print(prs.slides[1].shapes[7].table.cell(1,1).text'''
'''
for slide in prs.slides:
    for shape in slide.shapes:
        for table in shape.table:
            for r in table.rows:
                s = ""
                for c in r.cells:
                    s += c.text_frame.text + " | "
                print(s)

        if not shape.has_text_frame:
            continue

        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_runs.append(run.text)
                #print(run.text)
'''

#print(text_runs)