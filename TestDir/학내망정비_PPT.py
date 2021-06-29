from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

import glob

import win32com.client
from win32com.client import Dispatch

excel = win32com.client.Dispatch("Excel.Application")
excel.Visible = True
wb = excel.Workbooks.Open('D:\\Desktop\\7.xlsx')
ws = wb.ActiveSheet

arr = []

value1 = []
subVal = []
for i in range(2, 160):
    a = [ws.Cells(i, 1).Value, ws.Cells(i, 2).Value, ws.Cells(i, 3).Value, ws.Cells(i, 4).Value, ws.Cells(i, 5).Value, \
         ws.Cells(i, 6).Value, ws.Cells(i, 7).Value, ws.Cells(i, 8).Value, ws.Cells(i, 9).Value, ws.Cells(i, 10).Value]

    # [순위, 지역, 학교명, 담당자, ap1, ap2, poe1, 날짜]
    arr.append(a)


f_list = glob.glob("D:\\Desktop\\TEST10\\*\\*")
aa = f_list[0].split('\\')
for f_val in f_list:
    aa = f_val.split('\\')
    bb = aa[3].split('_')

    if aa[4].find("망구성도") != -1 and aa[4].find('.pptx') != -1:
        subVal = [f_val, arr[int(bb[0])-1][1], arr[int(bb[0])-1]]
        value1.append(subVal)

#print(value1)
#print(len(value1))
#print(len(f_list))

for item in value1:
    print(item[0])

    prs = Presentation(item[0])
    print((prs.slides[0].shapes[2].has_table and prs.slides[0].shapes[6].has_table))

    if prs.slides[0].shapes[2].has_table and prs.slides[0].shapes[6].has_table:
        #print(prs.slides[0].shapes[2].table.cell(5, 2).text)

        prs.slides[0].shapes[2].table.cell(5, 2).text = "기타(서비스망)"
        prs.slides[0].shapes[2].table.cell(5, 2).text_frame.paragraphs[0].font.name = '나눔고딕 Bold'
        prs.slides[0].shapes[2].table.cell(5, 2).text_frame.paragraphs[0].font.size = Pt(10)
        prs.slides[0].shapes[2].table.cell(5, 2).text_frame.paragraphs[0].font.bold = True
        prs.slides[0].shapes[2].table.cell(5, 2).text_frame.paragraphs[0].font.italic = None
        prs.slides[0].shapes[2].table.cell(5, 2).text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

        prs.slides[0].shapes[6].table.cell(4, 1).text = "기타(서비스망)"
        prs.slides[0].shapes[6].table.cell(4, 1).text_frame.paragraphs[0].font.name = '나눔고딕 Bold'
        prs.slides[0].shapes[6].table.cell(4, 1).text_frame.paragraphs[0].font.size = Pt(10)
        prs.slides[0].shapes[6].table.cell(4, 1).text_frame.paragraphs[0].font.bold = True
        prs.slides[0].shapes[6].table.cell(4, 1).text_frame.paragraphs[0].font.italic = None
        prs.slides[0].shapes[6].table.cell(4, 1).text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

        #print(item[1])
        #print(type(item[1]))
        #print(item[2][2])

        if item[1]:
            #print("OK")

            prs.slides[0].shapes[2].table.cell(2, 4).text = item[2][2]
            prs.slides[0].shapes[2].table.cell(3, 4).text = item[2][3]
            prs.slides[0].shapes[2].table.cell(4, 4).text = item[2][4]
            prs.slides[0].shapes[2].table.cell(5, 4).text = item[2][5]

            prs.slides[0].shapes[6].table.cell(1, 2).text = item[2][6]
            prs.slides[0].shapes[6].table.cell(2, 2).text = item[2][7]
            prs.slides[0].shapes[6].table.cell(3, 2).text = item[2][8]
            prs.slides[0].shapes[6].table.cell(4, 2).text = item[2][9]

            for i in [2, 3, 4, 5]:
                prs.slides[0].shapes[2].table.cell(i, 4).text_frame.paragraphs[0].font.name = '나눔고딕 Bold'
                prs.slides[0].shapes[2].table.cell(i, 4).text_frame.paragraphs[0].font.size = Pt(10)
                prs.slides[0].shapes[2].table.cell(i, 4).text_frame.paragraphs[0].font.bold = True
                prs.slides[0].shapes[2].table.cell(i, 4).text_frame.paragraphs[0].font.italic = None
                prs.slides[0].shapes[2].table.cell(i, 4).text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

                prs.slides[0].shapes[6].table.cell(i-1, 2).text_frame.paragraphs[0].font.name = '나눔고딕 Bold'
                prs.slides[0].shapes[6].table.cell(i-1, 2).text_frame.paragraphs[0].font.size = Pt(10)
                prs.slides[0].shapes[6].table.cell(i-1, 2).text_frame.paragraphs[0].font.bold = True
                prs.slides[0].shapes[6].table.cell(i-1, 2).text_frame.paragraphs[0].font.italic = None
                prs.slides[0].shapes[6].table.cell(i-1, 2).text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    prs.save(item[0])
