from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

prs = Presentation(".\\1.pptx")
main_slide = prs.slides[0]

for shape in prs.slides[0].shapes:
    if not shape.has_text_frame:
        continue
    if "기관명 :" in shape.text_frame.paragraphs[0].text:
        # print(shape.text_frame.paragraphs[0].text)
        shape.text_frame.paragraphs[0].text = "기관명 : 서울서울서울서울서울서울서울서울서울초등학교"
        for run in shape.text_frame.paragraphs[0].runs:
            run.font.size = Pt(16)
            run.font.name = '맑은 고딕'
            run.font.bold = True

for shape in prs.slides[0].shapes:
    if not shape.has_table:
        continue
    if "Port" in shape.table.cell(0, 0).text:
        # print(shape.table.cell(0, 0).text)
        shape.table.cell(2, 3).text = "20.20.20.20/24\n20.20.30.30/24"
        for paragraph in shape.table.cell(2, 3).text_frame.paragraphs:
            paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(12)
                run.font.name = '맑은 고딕'
                run.font.bold = True

        shape.table.cell(3, 3).text = "30.30.30.30/24"
        for paragraph in shape.table.cell(3, 3).text_frame.paragraphs:
            paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(12)
                run.font.name = '맑은 고딕'
                run.font.bold = True

        shape.table.cell(4, 3).text = "40.40.40.40/24"
        for paragraph in shape.table.cell(4, 3).text_frame.paragraphs:
            paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(12)
                run.font.name = '맑은 고딕'
                run.font.bold = True

        shape.table.cell(5, 3).text = "50.50.50.50/24"
        for paragraph in shape.table.cell(5, 3).text_frame.paragraphs:
            paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(12)
                run.font.name = '맑은 고딕'
                run.font.bold = True

prs.save(".\\1.pptx")
